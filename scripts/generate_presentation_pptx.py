#!/usr/bin/env python3
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt

INPUT_MD = os.path.join(os.path.dirname(__file__), '..', 'PRESENTATION_INSTRUCTIONS.md')
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'output')
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, 'Presentation_Instructions.pptx')


def read_md(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()


def split_slides(md_text):
    # Find slide sections that start with '**Slide'
    lines = md_text.splitlines()
    slides = []
    current = None
    for line in lines:
        if line.strip().startswith('**Slide'):
            if current:
                slides.append(current)
            # Start new section
            current = {'title': line.strip().lstrip('*').strip(), 'body_lines': []}
        else:
            if current is None:
                continue
            current['body_lines'].append(line)
    if current:
        slides.append(current)
    return slides


def extract_notes(body_lines):
    notes = []
    body = []
    capture = False
    for i, ln in enumerate(body_lines):
        stripped = ln.strip()
        if stripped.startswith('Speaker notes:') or stripped.startswith('- Speaker notes:'):
            capture = True
            continue
        if capture:
            if stripped == '' and len(notes) > 0:
                # stop on blank line after notes
                capture = False
                continue
            notes.append(ln)
        else:
            body.append(ln)
    # Clean up
    notes_text = '\n'.join([l.strip().lstrip('- ').strip() for l in notes if l.strip()!=''])
    body_text = '\n'.join([l.strip() for l in body if l.strip()!=''])
    return body_text, notes_text


def make_pptx(slides_sections):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = 'Privacy-Preserving Credit for the Unbanked'
    try:
        slide.placeholders[1].text = 'USDw Stablecoin + Zero-Knowledge Proofs\nPresentation Instructions'
    except Exception:
        pass

    # Create a slide for each section
    for s in slides_sections:
        title = s.get('title', '')
        # Clean title: remove leading ** and trailing timing info in parentheses
        title_clean = re.sub(r"\*+", '', title).strip()
        # Use em-dash or long dash normalization
        title_clean = title_clean.replace('—', '-')

        body_text, notes_text = extract_notes(s.get('body_lines', []))

        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_clean

        # Add body text as paragraphs
        tf = slide.shapes.placeholders[1].text_frame
        # If body_text empty, add a short hint
        if not body_text:
            p = tf.paragraphs[0]
            p.text = '(See speaker notes)'
        else:
            # Try splitting into bullet points by lines or by sentences
            for i, line in enumerate(body_text.split('\n')):
                if i == 0:
                    tf.text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0

        # Add speaker notes if present
        if notes_text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text

    # Ensure output dir
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


def main():
    print('Reading', INPUT_MD)
    md = read_md(INPUT_MD)
    slides = split_slides(md)
    if not slides:
        print('No slides found with pattern **Slide. Creating default slide from file.')
        slides = [{'title': 'Presentation Notes', 'body_lines': md.splitlines()}]
    out = make_pptx(slides)
    print('Saved PPTX to', out)


if __name__ == '__main__':
    main()
